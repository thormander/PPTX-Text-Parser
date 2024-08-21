import os
import re
import pandas as pd
from pptx import Presentation
from tqdm import tqdm

def contains_meaningful_content(text):
    """Check if the text contains any letters or numbers."""
    return bool(re.search(r'[a-zA-Z0-9]', text))

# Function to extract text from a shape
def extract_text_from_shape(shape):
    text = ""
    if hasattr(shape, "text_frame") and shape.text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text.strip() and contains_meaningful_content(run.text):
                    text += run.text + " "
    return text.strip()

# Function to extract text from a table
def extract_text_from_table(table):
    text = ""
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip() and contains_meaningful_content(run.text):
                        text += run.text + " "
    return text.strip()

# Recursive function to process all shapes in a slide
def process_shapes_recursive(shapes):
    text = ""
    for shape in shapes:
        if shape.has_text_frame:
            text += extract_text_from_shape(shape) + "\n"
        elif shape.has_table:
            text += extract_text_from_table(shape.table) + "\n"
        elif hasattr(shape, 'shapes'):  # Check if it's a group shape
            text += process_shapes_recursive(shape.shapes) + "\n"
    return text.strip()

# Process a single presentation and extract text
def process_presentation(input_file):
    print(f"Opening {input_file}")
    try:
        presentation = Presentation(input_file)
    except Exception as e:
        print(f"Error opening file {input_file}: {e}")
        return None

    text_data = []
    slide_count = len(presentation.slides)
    file_name = os.path.splitext(os.path.basename(input_file))[0]
    
    with tqdm(total=slide_count, desc="Extracting", unit="slide") as pbar:
        for slide_num, slide in enumerate(presentation.slides, start=1):
            slide_text = process_shapes_recursive(slide.shapes)
            if slide_text:
                slide_name_number = f"{file_name}_Slide{slide_num}"
                text_data.append({
                    "FileName_SlideNumber": slide_name_number,
                    "Slide Text": slide_text
                })
            pbar.update(1)
    
    return text_data

# Process all presentations in a folder
def process_folder(folder_path):
    all_text_data = []
    for filename in os.listdir(folder_path):
        if filename.endswith(".pptx"):
            file_path = os.path.join(folder_path, filename)
            text_data = process_presentation(file_path)
            if text_data:
                all_text_data.extend(text_data)
    return all_text_data

# Main function to parse arguments and initiate processing
def main():
    import argparse
    
    parser = argparse.ArgumentParser(description="Extract text from PowerPoint presentations.")
    parser.add_argument("input_path", nargs='?', help="Path to the input PowerPoint file or folder")
    args = parser.parse_args()

    if not args.input_path:
        parser.print_help()
        return
    
    # Determine output file path (in the same directory as the script)
    script_dir = os.path.dirname(os.path.abspath(__file__))
    output_file = os.path.join(script_dir, "extracted_slide_texts.csv")
    
    # Process individual file or folder
    if os.path.isdir(args.input_path):
        all_text_data = process_folder(args.input_path)
    else:
        all_text_data = process_presentation(args.input_path)
    
    if all_text_data:
        # Convert to DataFrame and save to CSV in the script's directory
        output_df = pd.DataFrame(all_text_data)
        output_df.to_csv(output_file, index=False)
        print(f"Extracted text saved to {output_file}")
    else:
        print("No text extracted.")

if __name__ == "__main__":
    main()
